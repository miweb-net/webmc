from pptx import Presentation
import sys


def process_file(file):

	prefix = file[:file.rfind('.')]

	p = Presentation(file)
	page = 0
	for slide in p.slides:
		page += 1
		img_num = 0
		for shape in slide.shapes:
			if shape.shape_type == 13:
				img_num += 1
				image = shape.image
				out_file = prefix + '_' + str(page) + '_' + str(img_num)
				out_file += '.png'
				print("   Writing:", out_file)
				with open(out_file, 'wb') as img_file:
					img_file.write(image.blob)

if __name__ == '__main__':

	# Get command line parameters or fail.
	if len(sys.argv) != 2:
		print("Usage:")
		print("    python " + sys.argv[0] + " <power point file in>")
		sys.exit(1)

	print("Extracting images from", sys.argv[1])
	process_file(sys.argv[1])