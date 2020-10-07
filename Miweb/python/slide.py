from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.text.text import Font
from pptx.util import Pt

from logger import Log
import io

log = Log()

class Slide():

	# Meta is a dictionary containing configuration elements for the
	# slide pages.
	def __init__(self, prefs, ratio=None):
		self.presentation = Presentation()

		self.prefs = prefs

		self.ccl_number = self.prefs.ccli.number
	
		# All slides in this program use blank layouts. So create one.
		self.blank_slide_layout = self.presentation.slide_layouts[6]

		# Slide width is 11" irrespective of aspect -- only height changes.
		self.presentation.slide_width = Inches(11)

		# Set the size of the presentation pages based on the configured
		if ratio == None:
			ratio = self.prefs.presentation.default
		if ratio == '16:9':
			self.presentation.slide_height = Inches(6.2)
			self.pg_h = 6.2
		else:
			self.presentation.slide_height = Inches(8.27)
			self.pg_h = 8.3
		log.config('Ratio set to:', ratio)

		self.pg_w = 11
	
	def addBlackSlide(self):
		self.slide = self.presentation.slides.add_slide(self.blank_slide_layout)
		self.slide.background.fill.solid()
		color = RGBColor.from_string('000000')
		self.slide.background.fill.fore_color.rgb = color

	def addWhiteSlide(self):
		self.slide = self.presentation.slides.add_slide(self.blank_slide_layout)
		self.slide.background.fill.solid()
		color = RGBColor.from_string('ffffff')
		self.slide.background.fill.fore_color.rgb = color

	def addTextSlide(self, slide, text, size=30):
		''' Write a text slide given the input text into a textbox element. The
		size argument is in points (e.g. 30 pt font)'''
		self.addBlackSlide()
		text_lines = text.split('\n')
		left, top = Inches(.3), Inches(.3)
		width, height = Inches(10.3), Inches(self.pg_h - 1.2)

		text_box = self.slide.shapes.add_textbox(left, top, width, height)
		text_frame = text_box.text_frame
		text_frame.word_wrap = True
		top_paragraph = text_frame.paragraphs[0]
		self.configureText(top_paragraph, text_lines[0], size, clr='ffffff',
					align=1, bold=True)
		
		for line in text_lines[1:]:
			next_paragraph = text_frame.add_paragraph()
			next_para = self.configureText(next_paragraph, line, size, clr='ffffff', 
			                            align=1, bold=True)

	def addScriptureTextReference(self, slide, passage):
		left, top, width, height = Inches(6), Inches(self.pg_h - 0.8), Inches(4.8), Inches(.5)
		text_box = self.slide.shapes.add_textbox(left, top, width, height)
		text_frame = text_box.text_frame
		text_frame.word_wrap = True
		top_paragraph = text_frame.paragraphs[0]
		self.configureText(top_paragraph, passage, 25, bold=True, italic=True,
									clr='ffffff', align=3)
		
	def configureText(self, para, text, size, align=2, clr='000000', bold=False,
					italic=False):
		para.alignment = align
		run = para.add_run()
		run.text = str(text)
		
		font = run.font
		font.name = 'Liberation Serif'
		font.size = Pt(int(size))
		font.bold = bold
		font.italic = italic
		font.fill.solid()
		color = RGBColor.from_string(clr)
		font.fill.fore_color.rgb = color
		
	def addMusicSlide(self, path, song, add_meta_data):
		''' Creates an image from the file designated by path.'''

		self.addWhiteSlide()
		left, top = Inches(0.3), Inches(0.4)
		pic = self.slide.shapes.add_picture(path, left, top, width=Inches(10.4), height=Inches(self.pg_h - 1))

		# Hymn Number placement
		left,top,width,height = Inches(8.6), Inches(0), Inches(2), Inches(.7)
		frame = self.createParagraph(left, top, width, height)
		self.configureText(frame, song['number'], 26, bold=True, align=3)

		# Title at top-left of page
		left,top,width,height = Inches(0.5), Inches(0), Inches(8.4), Inches(0.7)
		frame = self.createParagraph(left, top, width, height)
		self.configureText(frame, song['title'], 26, bold=True, align=1)
		
		# The first slide must also include the author(s), date(s), copyright
		# and CCLI#. If page_id is the first page_id in the song's list of pages,
		# add the remaining meta-data to the page.
		if add_meta_data:
			self.addMusicMeta(song)
			self.addHiddenMetaRecord(self.slide, song)
		
	def addMusicMeta(self, song):
		# Copyright at center of page bottom
		if len(song['copyright']) and 'Public Domain' not in song['copyright']:
			left, top = Inches(0.5), Inches(self.pg_h - 0.35)
			width, height = Inches(8.5), Inches(0.3)
			frame = self.createParagraph(left, top, width, height)
			copyright = song['copyright'] + ', CCL No. ' + self.ccl_number
			self.configureText(frame, copyright, 13, bold=False,
							   align=1, clr='555555')
		# Author (lyricist) below song number (at top-left of page).
		if len(song['lyricist']):
			author_txt = 'Lyrics: ' + song['lyricist']
			date = str(song['ldate'])
			if len(date):
				author_txt += ', ' + date
		# Author text is located in upper-left corner
		left, top, width, height = Inches(0.5), Inches(0.5), \
								   Inches(4.0), Inches(0.3)
		frame = self.createParagraph(left, top, width, height)
		self.configureText(frame, author_txt, 13, bold=False, align=1)

		# Composer (arranger) below song title (at top-right of page).
		if len(song['composer']):
			composer_txt = 'Music: ' + song['composer']
			date = str(song['cdate'])
			if len(date):
				composer_txt += ', ' + date
		left,top,width,height = Inches(6.55), Inches(0.5), Inches(4.0), Inches(0.3)
		frame = self.createParagraph(left, top, width, height)
		self.configureText(frame, composer_txt, 13, bold=False, align=3)

		# And finally, the MIMC logo
		path = '../image/desktop_icon.png'
		left, top = Inches(9.4), Inches(self.pg_h - .6)
		pic = self.slide.shapes.add_picture(path, left, top, width=Inches(.5))

		
	def createParagraph(self, left, top, width, height):
		text_box = self.slide.shapes.add_textbox(left, top, width, height)
		text_frame = text_box.text_frame
		text_frame.word_wrap = True
		return text_frame.paragraphs[0]

	def addHiddenMetaRecord(self, slide, song):
		left, top, width, height = Inches(10), Inches(9), Inches(.5), Inches(.3)
		frame = self.createParagraph(left, top, width, height)

		embed_data = str(song) + str(self.prefs)
		self.configureText(frame, str(self.prefs), 5, clr='c5c5c5', bold=False, align=1)

	def addMedleySlide(self, song):
		''' Create a text slide with only hymn meta-data included in the medley.
			Format is Book, Number, Title, Lyricist, (Date), Composer, (Date),
			Copyright for each hymn.
		'''
		self.addBlackSlide()
		color = '88aaff'

		# CCLI Number
		left = Inches(8)
		frame = self.createParagraph(left, Inches(.4), Inches(5), Inches(.3))
		ccli = "CCLI: " + self.prefs.ccli.number
		self.configureText(frame, ccli, 16, clr=color, bold=False, align=1)

		frame = self.createParagraph(Inches(2.7), Inches(.4), Inches(5), Inches(.3))
		text = "A Medley Presentation"
		self.configureText(frame, text, 24, clr=color, bold=True, align=1)

		left = Inches(1.6)
		top = Inches(1.1)
		#for book in meta['BOOK'].split('~'):
		#	frame = self.createParagraph(left, top, Inches(5), Inches(.3))
		#	self.configureText(frame, book, 20, clr=color, bold=False, align=1)
		#	top += Inches(2.6)

		top = Inches(1.5)
		for number in song['number']:
			frame = self.createParagraph(left, top, Inches(.8), Inches(.3))
			self.configureText(frame, number, 20, clr=color, bold=True, align=1)
			top += Inches(2.6)

		top = Inches(1.5)
		for title in song['titles']:
			frame = self.createParagraph(left + Inches(.5), top, Inches(8), Inches(.3))
			self.configureText(frame, title, 20, clr=color, bold=True, align=1)
			top += Inches(2.6)

		top = Inches(1.9)
		for ldate,lyricist in zip(song['ldate'], song['lyricist']):
			text = "Lyrics: " + lyricist
			if ldate != "Unknown":
				text += ", " + ldate
			frame = self.createParagraph(left, top, Inches(8), Inches(.3))
			self.configureText(frame, text, 20, clr=color, bold=False, align=1)
			top += Inches(2.6)

		top = Inches(2.3)
		for cdate,composer in zip(song['cdate'], song['composer']):
			text = "Composer: " + composer
			if cdate != "Unknown":
				text += ", " + cdate
			frame = self.createParagraph(left, top, Inches(8), Inches(.3))
			self.configureText(frame, text, 20, clr=color, bold=False, align=1)
			top += Inches(2.6)

		top = Inches(2.7)
		width = Inches(9)
		for copyright in song['copyright']:
			frame = self.createParagraph(left, top, width, Inches(.3))
			self.configureText(frame, copyright, 14, clr=color, bold=False, align=1)
			top += Inches(2.6)

		# And finally, the MIMC logo
		path = '../image/desktop_icon.png'
		left, top = Inches(9), Inches(6.9)
		pic = self.slide.shapes.add_picture(path, left, top, width=Inches(.5))

	def saveFile(self, filename):
		'''Save the file as filename and return filename to the caller.'''
		self.presentation.save(filename)
		return filename
		
